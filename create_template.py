"""
Create a PowerPoint template with placeholders for the marketing tool.

This script creates a PowerPoint template with placeholders (x1-x15) that will be
replaced with actual data when generating reports.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from pathlib import Path

def create_template():
    """Create a PowerPoint template with placeholders."""
    # Create a new presentation
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Marketing Campaign Report for x1"
    subtitle.text = "Reporting Period: x2"
    
    # Slide 2: Overview
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Campaign Overview"
    
    tf = content.text_frame
    tf.text = "Key Performance Indicators:"
    
    p = tf.add_paragraph()
    p.text = "Total Impressions: x3"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total Clicks: x4"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total Conversions: x5"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total Spend: x6"
    p.level = 1
    
    # Slide 3: Metrics
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Campaign Metrics"
    
    tf = content.text_frame
    tf.text = "Performance Metrics:"
    
    p = tf.add_paragraph()
    p.text = "Click-Through Rate (CTR): x7"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Conversion Rate: x8"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cost Per Click (CPC): x9"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cost Per Conversion: x10"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Return on Investment (ROI): x11"
    p.level = 1
    
    # Slide 4: Trends
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Performance Trends"
    
    tf = content.text_frame
    tf.text = "Week-over-Week Change: x12"
    
    p = tf.add_paragraph()
    p.text = "Month-over-Month Change: x13"
    p.level = 1
    
    # Slide 5: Campaign Details
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Campaign Details"
    
    tf = content.text_frame
    tf.text = "Active Campaigns: x14"
    
    p = tf.add_paragraph()
    p.text = "Total Budget: x15"
    p.level = 1
    
    # Save the template
    template_dir = Path("app/templates")
    template_dir.mkdir(exist_ok=True)
    
    template_path = template_dir / "report_template.pptx"
    prs.save(template_path)
    
    print(f"Template created at: {template_path}")
    return template_path

if __name__ == "__main__":
    template_path = create_template()
    print(f"PowerPoint template created successfully at {template_path}")
    print("The template contains the following placeholders:")
    print("x1: Client name")
    print("x2: Date range")
    print("x3: Total impressions")
    print("x4: Total clicks")
    print("x5: Total conversions")
    print("x6: Total spend")
    print("x7: CTR")
    print("x8: Conversion rate")
    print("x9: Cost per click")
    print("x10: Cost per conversion")
    print("x11: ROI")
    print("x12: Week-over-week change")
    print("x13: Month-over-month change")
    print("x14: Number of campaigns")
    print("x15: Total budget")
