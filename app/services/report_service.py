from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any, List
import pandas as pd
from datetime import datetime
import logging
import os

logger = logging.getLogger(__name__)

class ReportService:
    def __init__(self):
        self.template_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates")
        self.template_path = os.path.join(self.template_dir, "report_template.pptx")
        if not os.path.exists(self.template_dir):
            os.makedirs(self.template_dir)
            logger.warning(f"Template directory created at {self.template_dir}")
        if not os.path.exists(self.template_path):
            logger.warning(f"Template file not found at {self.template_path}")

    def generate_report(
        self,
        template_path: str,
        campaign_data: Dict[str, Any],
        metrics_data: Dict[str, Any],
        client_name: str,
        date_range: Dict[str, datetime]
    ) -> str:
        """Generate a PowerPoint report using the provided template and data"""
        try:
            prs = Presentation(template_path)
            
            # Title slide
            title_slide = prs.slides[0]
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = f"Campaign Report - {client_name}"
            subtitle.text = f"Period: {date_range['start'].strftime('%B %d, %Y')} - {date_range['end'].strftime('%B %d, %Y')}"

            # Campaign Overview slide
            overview_slide = prs.add_slide(prs.slide_layouts[1])
            title = overview_slide.shapes.title
            content = overview_slide.placeholders[1]
            
            title.text = "Campaign Overview"
            content.text = self._format_campaign_overview(campaign_data)

            # Metrics slide
            metrics_slide = prs.add_slide(prs.slide_layouts[2])
            title = metrics_slide.shapes.title
            content = metrics_slide.placeholders[1]
            
            title.text = "Campaign Metrics"
            content.text = self._format_metrics(metrics_data)

            # Budget slide
            budget_slide = prs.add_slide(prs.slide_layouts[2])
            title = budget_slide.shapes.title
            content = budget_slide.placeholders[1]
            
            title.text = "Budget Utilization"
            content.text = self._format_budget_data(campaign_data)

            # Save the presentation
            output_path = f"reports/{client_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(output_path)
            
            return output_path

        except Exception as e:
            logger.error(f"Error generating report: {str(e)}")
            raise

    def _format_campaign_overview(self, campaign_data: Dict[str, Any]) -> str:
        """Format campaign overview data for the report"""
        return f"""
Campaign Name: {campaign_data['name']}
Platform: {campaign_data['platform']}
Status: {campaign_data['status']}
Start Date: {campaign_data['start_date']}
End Date: {campaign_data['end_date'] or 'Ongoing'}
Total Budget: ${campaign_data['total_budget']:,.2f}
Spent Budget: ${campaign_data['spent_budget']:,.2f}
"""

    def _format_metrics(self, metrics_data: Dict[str, Any]) -> str:
        """Format metrics data for the report"""
        # Handle both LinkedIn and Rollworks metrics
        if "impressions" in metrics_data:
            # LinkedIn format
            return f"""
Impressions: {metrics_data['impressions']:,}
Clicks: {metrics_data['clicks']:,}
CTR: {metrics_data['ctr']:.2%}
Conversions: {metrics_data['conversions']:,}
Conversion Rate: {metrics_data['conversion_rate']:.2%}
Cost per Click: ${metrics_data['cpc']:.2f}
Cost per Conversion: ${metrics_data['cost_per_conversion']:.2f}
"""
        else:
            # Rollworks format
            return f"""
Impressions: {metrics_data.get('impressions', 0):,}
Clicks: {metrics_data.get('clicks', 0):,}
CTR: {metrics_data.get('ctr', 0):.2%}
Conversions: {metrics_data.get('conversions', 0):,}
Conversion Rate: {metrics_data.get('conversion_rate', 0):.2%}
Cost per Click: ${metrics_data.get('cpc', 0):.2f}
Cost per Conversion: ${metrics_data.get('cost_per_conversion', 0):.2f}
View-through Conversions: {metrics_data.get('view_through_conversions', 0):,}
Click-through Conversions: {metrics_data.get('click_through_conversions', 0):,}
"""

    def _format_budget_data(self, campaign_data: Dict[str, Any]) -> str:
        """Format budget data for the report"""
        spent_percentage = (campaign_data['spent_budget'] / campaign_data['total_budget']) * 100
        return f"""
Total Budget: ${campaign_data['total_budget']:,.2f}
Spent Budget: ${campaign_data['spent_budget']:,.2f}
Budget Utilization: {spent_percentage:.1f}%
Remaining Budget: ${campaign_data['total_budget'] - campaign_data['spent_budget']:,.2f}
""" 