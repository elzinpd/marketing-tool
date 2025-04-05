import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional
import random
from datetime import timedelta

from fastapi import APIRouter, Depends, Header, HTTPException, Request, status
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session

from app.core.config import settings
from app.core.database import get_db
from app.core.auth import get_current_active_user
from app.db.models import User, Client
from app.services.linkedin_service import LinkedInService
from pptx import Presentation

router = APIRouter()
logger = logging.getLogger(__name__)

# Create reports directory if it doesn't exist
os.makedirs(settings.REPORTS_DIRECTORY, exist_ok=True)
reports_dir_abs = os.path.abspath(settings.REPORTS_DIRECTORY)
logger.info(f"Reports directory: {reports_dir_abs}")

@router.get("/test", response_model=Dict[str, str])
async def test_reports_endpoint():
    """
    Test endpoint to verify the reports API is working.
    """
    return {"status": "ok", "message": "Reports API is operational"}

@router.get("/export-client/{client_id}")
async def export_client_report(
    request: Request,
    client_id: int,
    start_date: str,
    end_date: str,
    authorization: str = Header(None),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Export client report as PowerPoint file
    """
    logger.info(f"Processing report export for client_id: {client_id}, dates: {start_date} to {end_date}")
    
    # Double check authorization
    if not authorization:
        logger.warning("Authorization header missing in report export request")
        raise HTTPException(
            status_code=401,
            detail="Authorization header is required",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    # Verify token format is correct (this is just a logging check, actual auth happens via get_current_active_user)
    if not authorization.startswith("Bearer "):
        logger.warning("Authorization header does not contain Bearer token")
        raise HTTPException(
            status_code=401, 
            detail="Invalid authorization format. Bearer token required",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    # Check client access
    client = get_client_by_id(db, client_id)
    if not client:
        logger.error(f"Client with id {client_id} not found")
        raise HTTPException(status_code=404, detail="Client not found")
    
    # Check user access to client
    user_has_access = check_user_client_access(db, current_user.id, client_id)
    if not user_has_access:
        logger.warning(f"User {current_user.email} attempted to access client {client_id} without permission")
        raise HTTPException(status_code=403, detail="Not authorized to access this client")
    
    try:
        # Define the report output path
        report_dir = Path(settings.REPORTS_DIRECTORY)
        report_dir.mkdir(exist_ok=True)
        
        filename = f"client_{client_id}_report_{start_date}_to_{end_date}.pptx"
        output_path = report_dir / filename
        
        # Generate PowerPoint
        create_powerpoint_report(
            output_path=output_path,
            client_name=client.name,
            start_date=start_date,
            end_date=end_date
        )
        
        # Return file as response
        return FileResponse(
            path=output_path,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        logger.exception(f"Error generating PowerPoint report: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate report: {str(e)}")

def create_powerpoint_report(output_path, client_name, start_date, end_date):
    """
    Create a PowerPoint report for a client with the given data using the
    template with labeled placeholders x1-x7.
    """
    try:
        # Parse date strings to datetime objects
        start_datetime = datetime.strptime(start_date, "%Y-%m-%d")
        end_datetime = datetime.strptime(end_date, "%Y-%m-%d")
        
        # Get mock campaign data for this client - use name as identifier
        from app.api.api_v1.endpoints.linkedin import get_linkedin_campaigns
        campaigns = get_linkedin_campaigns(client_name=client_name)
        
        # Calculate total metrics
        total_impressions = sum(campaign.get("metrics", {}).get("impressions", 0) for campaign in campaigns)
        total_clicks = sum(campaign.get("metrics", {}).get("clicks", 0) for campaign in campaigns)
        total_conversions = sum(campaign.get("metrics", {}).get("conversions", 0) for campaign in campaigns)
        total_spend = sum(campaign.get("metrics", {}).get("spend", 0) for campaign in campaigns)
        
        # Calculate derived metrics
        ctr = (total_clicks / total_impressions) if total_impressions > 0 else 0
        conversion_rate = (total_conversions / total_clicks) if total_clicks > 0 else 0
        cpc = (total_spend / total_clicks) if total_clicks > 0 else 0
        cost_per_conversion = (total_spend / total_conversions) if total_conversions > 0 else 0
        roi = ((total_conversions * 100) - total_spend) / total_spend if total_spend > 0 else 0
        
        # Calculate week-over-week and month-over-month changes (mock data)
        wow_change = random.uniform(-0.15, 0.25)  # -15% to +25% change
        mom_change = random.uniform(-0.1, 0.3)    # -10% to +30% change
        
        # Calculate campaign performance over time (for trend analysis)
        date_range = (end_datetime - start_datetime).days
        performance_data = {
            "dates": [],
            "impressions": [],
            "clicks": [],
            "conversions": [],
            "spend": []
        }
        
        # Generate daily data points for the time series
        random.seed(hash(client_name) % 1000)
        for i in range(date_range + 1):
            current_date = start_datetime + timedelta(days=i)
            performance_data["dates"].append(current_date.strftime("%Y-%m-%d"))
            
            # Generate values with slight day-to-day variations for realistic trends
            daily_factor = 1 + (random.random() * 0.2 - 0.1)  # -10% to +10% daily variation
            
            daily_impressions = int((total_impressions / date_range) * daily_factor)
            daily_clicks = int((total_clicks / date_range) * daily_factor)
            daily_conversions = int((total_conversions / date_range) * daily_factor)
            daily_spend = (total_spend / date_range) * daily_factor
            
            performance_data["impressions"].append(daily_impressions)
            performance_data["clicks"].append(daily_clicks)
            performance_data["conversions"].append(daily_conversions)
            performance_data["spend"].append(daily_spend)
        
        # Load PowerPoint template from templates directory
        template_path = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), 
                                    "templates", "report_template.pptx")
        
        # Check if template exists
        if not os.path.exists(template_path):
            logger.error(f"Template file not found at: {template_path}")
            raise Exception(f"Report template not found at: {template_path}")
            
        logger.info(f"Using template from: {template_path}")
        prs = Presentation(template_path)
        
        # Dictionary mapping placeholder labels to their values
        placeholder_values = {
            "x1": client_name,                                             # Client name
            "x2": f"{start_date} to {end_date}",                           # Date range
            "x3": f"{total_impressions:,}",                                # Total impressions
            "x4": f"{total_clicks:,}",                                     # Total clicks
            "x5": f"{total_conversions:,}",                                # Total conversions
            "x6": f"${total_spend:,.2f}",                                  # Total spend
            "x7": f"{ctr:.2%}",                                            # CTR
            "x8": f"{conversion_rate:.2%}",                                # Conversion rate
            "x9": f"${cpc:.2f}",                                           # Cost per click
            "x10": f"${cost_per_conversion:.2f}",                          # Cost per conversion
            "x11": f"{roi:.1%}",                                           # ROI
            "x12": f"{wow_change:+.1%}",                                   # Week-over-week change
            "x13": f"{mom_change:+.1%}",                                   # Month-over-month change
            "x14": str(len(campaigns)),                                    # Number of campaigns
            "x15": f"${sum(c.get('budget', 0) for c in campaigns):,.2f}"   # Total budget
        }
        
        # Process each slide in the presentation
        for slide in prs.slides:
            # Replace placeholders in text boxes and shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_frame = shape.text_frame
                    
                    # Process each paragraph in the text frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Replace placeholder texts with actual values
                            for placeholder, value in placeholder_values.items():
                                if placeholder in run.text:
                                    run.text = run.text.replace(placeholder, value)
        
        # Add a campaign details slide if not already in template
        has_campaign_details = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and "Campaign Details" in shape.text:
                    has_campaign_details = True
                    break
        
        if not has_campaign_details:
            # Add a campaign details slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Campaign Details"
            
            # Add text to the slide
            tf = body.text_frame
            tf.text = "Active Campaigns"
            
            # Add campaign details
            for campaign in campaigns[:5]:  # Limit to first 5 campaigns
                p = tf.add_paragraph()
                p.text = f"{campaign.get('name')}: ${campaign.get('metrics', {}).get('spend', 0):,.2f}"
                p.level = 1
        
        # Save the presentation
        prs.save(output_path)
        logger.info(f"PowerPoint report created successfully at {output_path}")
        return str(output_path)
    except Exception as e:
        logger.exception(f"Error creating PowerPoint report: {str(e)}")
        raise

def get_client_by_id(db: Session, client_id: int):
    """Get a client by ID from the database."""
    return db.query(Client).filter(Client.id == client_id).first()

def check_user_client_access(db: Session, user_id: int, client_id: int):
    """Check if a user has access to a client. Admins have access to all clients."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        return False
    
    # Admins can access all clients
    if user.role == "admin":
        return True
    
    # Check client assignments for non-admin users
    user_client_ids = [c.id for c in user.clients]
    return client_id in user_client_ids 